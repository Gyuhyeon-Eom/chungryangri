"""발표자료(PPTX) 생성 — DESIGN-airtable.md 지침 적용.

지침의 핵심을 슬라이드로 옮긴 방식
  - 슬라이드 하나가 지침의 '에디토리얼 밴드' 하나에 대응한다.
  - 표면 모드를 교대시킨다: 흰 캔버스가 기본이고, 시그니처 카드(coral/forest/dark/cream)가
    2~3장마다 들어가 리듬을 만든다. 지침의 "signature cards punctuate every two or three screens".
  - 디스플레이 타입은 볼드로 가지 않는다(weight 400~500). 강조는 크기와 색 대비로 준다.
  - 그림자·그라디언트를 쓰지 않는다. 깊이는 흰 캔버스와 시그니처 면의 대비로만 만든다.
  - 96px 섹션 리듬 → 슬라이드 여백 2.2cm 로 환산.

    uv run python src/build_pptx.py
"""

import subprocess
import tempfile
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

import charts as C
from config import DATA_PROCESSED, ROOT

C.use_theme("airtable")

OUT = ROOT / "output" / "청량리_전통시장_상권분석_발표자료.pptx"
TMP = Path(tempfile.mkdtemp(prefix="deck_"))

# --- 색 토큰 ---
# 지침의 구조(흰 캔버스 + 시그니처 카드 리듬, 볼드 없는 디스플레이)는 그대로 두고
# 색만 청량리종합시장 성격에 맞춰 바꿨다. 주력 품목이 밤·견과류·곡류·먹거리라
# 차가운 네이비 대신 따뜻한 흙색 계열을 쓴다.
INK = RGBColor(0x2A, 0x24, 0x1E)          # 짙은 갈색빛 먹색 — 표지·본문 기준색
BODY = RGBColor(0x44, 0x3D, 0x35)
MUTED = RGBColor(0x6B, 0x62, 0x57)
CANVAS = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xFA, 0xF7, 0xF2)         # 곡물자루 같은 미색
HAIRLINE = RGBColor(0xDD, 0xD6, 0xCB)
CORAL = RGBColor(0x8C, 0x3A, 0x1E)        # 대추·밤 껍질색
FOREST = RGBColor(0x2D, 0x47, 0x39)       # 채소 진녹색
CREAM = RGBColor(0xF2, 0xE8, 0xD5)        # 한지·포대 미색
ON_DARK = RGBColor(0xFF, 0xFF, 0xFF)

FONT_KR = "Apple SD Gothic Neo"   # 지침: macOS에서는 system-ui 로 충분
MARGIN = Cm(2.2)                   # 96px 섹션 리듬 환산
W, H = Cm(33.87), Cm(19.05)        # 16:9


def set_run(run, size, color, bold=False, font=FONT_KR):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold          # 지침: 디스플레이는 볼드 금지
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", font)


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.paragraphs[0].alignment = align
    return tf


def add_text(tf, text, size, color, bold=False, space_after=0, line=1.3, first=False, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after)
    p.line_spacing = line
    if align is not None:
        p.alignment = align
    set_run(p.add_run(), size, color, bold)
    p.runs[0].text = text
    return p


def blank(prs, fill=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if fill is not None:
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill
        bg.line.fill.background()
        bg.shadow.inherit = False
    return s


def eyebrow(slide, text, color=MUTED):
    """지침의 caption 역할. 대문자·자간 넓힘으로 섹션을 표시한다."""
    tf = textbox(slide, MARGIN, Cm(1.5), W - MARGIN * 2, Cm(0.8))
    add_text(tf, text, 11, color, bold=True, first=True)


def title(slide, text, color=INK, y=Cm(2.4), size=30):
    """디스플레이 타입. 지침에 따라 볼드를 쓰지 않는다."""
    tf = textbox(slide, MARGIN, y, W - MARGIN * 2, Cm(3.0))
    for i, line in enumerate(text.split("\n")):
        add_text(tf, line, size, color, bold=False, line=1.2, first=(i == 0))


def rule(slide, y, color=INK, w=Cm(4.6)):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y, w, Pt(2.5))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False


def chart(slide, svg, name, x, y, w):
    png = TMP / f"{name}.png"
    (TMP / f"{name}.svg").write_text(svg, encoding="utf-8")
    subprocess.run(["rsvg-convert", "-w", "2200", "-b", "white",
                    "-o", str(png), str(TMP / f'{name}.svg')], check=True)
    slide.shapes.add_picture(str(png), x, y, width=w)


def table(slide, headers, rows, x, y, w, col_w, font=10.5, dark=False):
    n_r, n_c = len(rows) + 1, len(headers)
    gt = slide.shapes.add_table(n_r, n_c, x, y, w, Cm(0.9 * n_r)).table
    for i, cw in enumerate(col_w):
        gt.columns[i].width = Cm(cw)
    for i, h in enumerate(headers):
        c = gt.cell(0, i)
        c.fill.solid()
        c.fill.fore_color.rgb = INK if not dark else CREAM
        c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_run(c.text_frame.paragraphs[0].add_run(), font, ON_DARK if not dark else INK, True)
        c.text_frame.paragraphs[0].runs[0].text = h
    for r, row in enumerate(rows, 1):
        for i, v in enumerate(row):
            c = gt.cell(r, i)
            c.fill.solid()
            c.fill.fore_color.rgb = CANVAS if r % 2 else SOFT
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if i else PP_ALIGN.LEFT
            set_run(p.add_run(), font, BODY)
            p.runs[0].text = str(v)
    return gt


def signature(prs, fill, eyebrow_text, headline, body, stat=None, stat_label=None):
    """시그니처 카드 슬라이드. 지침의 voltage moment — 2~3장마다 하나씩."""
    on = ON_DARK if fill != CREAM else INK
    sub = RGBColor(0xE0, 0xE2, 0xE6) if fill != CREAM else BODY
    s = blank(prs, fill)
    tf = textbox(s, MARGIN, Cm(1.6), W - MARGIN * 2, Cm(0.8))
    add_text(tf, eyebrow_text, 11, sub, bold=True, first=True)
    rule(s, Cm(2.7), on)

    tf = textbox(s, MARGIN, Cm(3.6), Cm(18.5), Cm(6))
    for i, line in enumerate(headline.split("\n")):
        add_text(tf, line, 30, on, bold=False, line=1.25, first=(i == 0))

    tf = textbox(s, MARGIN, Cm(11.4), Cm(17.5), Cm(5))
    for i, line in enumerate(body):
        add_text(tf, line, 13, sub, line=1.5, space_after=6, first=(i == 0))

    if stat:
        tf = textbox(s, Cm(22.5), Cm(5.4), Cm(9.2), Cm(5), align=PP_ALIGN.RIGHT)
        add_text(tf, stat, 62, on, bold=False, line=1.0, first=True, align=PP_ALIGN.RIGHT)
        add_text(tf, stat_label, 12, sub, line=1.4, space_after=0, align=PP_ALIGN.RIGHT)
    return s


# ---------------------------------------------------------------- 본문
def build():
    bm = pd.read_csv(DATA_PROCESSED / "seoul_benchmark.csv")
    ss = pd.read_csv(DATA_PROCESSED / "shift_share.csv", index_col=0)
    dt = pd.read_csv(DATA_PROCESSED / "ticket_decomp.csv", index_col=0)
    el = pd.read_csv(DATA_PROCESSED / "elasticity_final.csv")
    clus = pd.read_csv(DATA_PROCESSED / "market_clusters.csv", index_col=0)
    dow = pd.read_csv(DATA_PROCESSED / "sales_dow.csv", index_col=0)
    ash = pd.read_csv(DATA_PROCESSED / "pop_age_shift.csv", index_col=0)
    dq = pd.read_csv(DATA_PROCESSED / "vacancy_data_quality.csv")

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1 — 표지 (dark)
    s = blank(prs, INK)
    tf = textbox(s, MARGIN, Cm(4.2), Cm(6), Cm(0.8))
    add_text(tf, "MARKET INTELLIGENCE REPORT", 11, RGBColor(0x9A, 0xA0, 0xAA), bold=True, first=True)
    rule(s, Cm(5.5), ON_DARK)
    tf = textbox(s, MARGIN, Cm(6.6), Cm(24), Cm(5))
    add_text(tf, "청량리 전통시장 클러스터", 40, ON_DARK, line=1.22, first=True)
    add_text(tf, "상권분석 및 Pain Point 진단", 40, ON_DARK, line=1.22)
    tf = textbox(s, MARGIN, Cm(13.0), Cm(22), Cm(2.4))
    add_text(tf, "서울 동대문구 제기동·청량리동·용신동 일대 9개 전통시장", 13,
             RGBColor(0xC0, 0xC4, 0xCA), line=1.5, first=True)
    add_text(tf, "2026. 8", 13, RGBColor(0x9A, 0xA0, 0xAA), line=1.5)

    # 2 — 분석 개요 (white)
    s = blank(prs, CANVAS)
    eyebrow(s, "01  분석 설계")
    title(s, "네 갈래 자료를 교차 검증했다")
    rule(s, Cm(4.6))
    table(s, ["자료", "단위", "기간", "용도"], [
        ["서울시 상권분석 추정매출", "상권 × 업종 × 분기", "2021~2025", "매출·시간대·요일·연령"],
        ["서울시 생활인구", "행정동 × 시간 × 연령", "2017~2026", "배후인구 변화"],
        ["소상공인시장진흥공단", "시장 × 연", "2006~2023", "점포·공실"],
        ["국세청 상가업소정보", "점포 좌표 × 업종", "2026.06", "공간 검증"],
    ], MARGIN, Cm(5.7), W - MARGIN * 2, [8.0, 8.0, 5.5, 8.0])
    tf = textbox(s, MARGIN, Cm(13.4), W - MARGIN * 2, Cm(3))
    add_text(tf, "생활인구 원시 레코드 5,430만 행 · 원본 138GB · 전국 1,812개 시장 교차검증", 13, BODY, first=True)
    add_text(tf, "변이-할당 분석 · 고정효과 패널회귀 · 계층적 군집분석 · 매출 요인분해 · 공간 조인",
             13, MUTED, space_after=0)

    # 3 — 핵심 발견 (cream)
    s = blank(prs, CREAM)
    eyebrow(s, "핵심 발견", BODY)
    title(s, "문제는 쇠퇴가 아니라 성장의 내용이다")
    rule(s, Cm(4.6))
    tf = textbox(s, MARGIN, Cm(7.0), Cm(14.5), Cm(10))
    for i, t in enumerate([
        "01   서울 전통시장 283개 중 8개 시장이 상위 20% 이내",
        "02   9개 중 8개에서 객단가 하락 — 물량으로 버티는 구조",
        "03   서울약령시장만 매출 감소, 원인은 업종이 아닌 입지",
        "04   배후인구는 줄지만 매출과의 연결은 확인되지 않음",
        "05   집객 성과가 고객 구성 변화로 전환되지 않음",
        "06   현행 공실 통계는 정책 근거로 사용 불가",
    ]):
        add_text(tf, t, 15, INK, line=1.45, space_after=14, first=(i == 0))
    tf = textbox(s, Cm(19.0), Cm(7.4), Cm(12.6), Cm(8), align=PP_ALIGN.RIGHT)
    add_text(tf, "5,444", 66, INK, line=1.0, first=True)
    add_text(tf, "인구-매출 관계 검정에 사용한 관측치", 12, BODY, line=1.4, space_after=18)
    add_text(tf, "283", 66, INK, line=1.0)
    add_text(tf, "비교 대상 서울 전통시장 상권", 12, BODY, line=1.4)

    # 4 — 서울 대비 위치 (white)
    s = blank(prs, CANVAS)
    eyebrow(s, "02  상권 위치")
    title(s, "청량리는 서울 전통시장 상위권이다")
    rule(s, Cm(4.6))
    chart(s, C.hbar([(r["name"], float(r["백분위"])) for _, r in
                     bm[bm["청량리"]].sort_values("백분위").iterrows()],
                    unit="%", label="서울 내 백분위", highlight={"서울약령시장"}), "bench", MARGIN, Cm(5.7), Cm(19.4))
    tf = textbox(s, Cm(21.6), Cm(6.0), Cm(10), Cm(10))
    add_text(tf, "서울 전통시장 283개의 2021~2025년 매출 성장률 분포에서 청량리 9개의 위치.",
             13, BODY, line=1.5, space_after=14, first=True)
    add_text(tf, "서울 중위 성장률 4.9%", 12, MUTED, line=1.5, space_after=4)
    add_text(tf, "청량리수산시장 상위 1%", 13, INK, line=1.5, space_after=4)
    add_text(tf, "동서시장 상위 3%", 13, INK, line=1.5, space_after=14)
    add_text(tf, "서울약령시장만 하위 25% 구간", 13, CORAL, line=1.5)

    # 5 — 매출 구조 (white)
    s = blank(prs, CANVAS)
    eyebrow(s, "03  성장의 구성")
    title(s, "거래는 늘고, 단가는 떨어졌다")
    rule(s, Cm(4.6))
    chart(s, C.scatter([(m, float(r["건수증가%"]), float(r["객단가증가%"]),
                         "#aa2d00" if r["객단가증가%"] < -30 else "#181d26")
                        for m, r in dt.iterrows()],
                       xlab="거래 건수 증가율 (%)", ylab="객단가 증가율 (%)",
                       quadrant=(0, 0), height=250, label="건수와 객단가"),
          "ticket", MARGIN, Cm(5.7), Cm(19.4))
    tf = textbox(s, Cm(21.6), Cm(6.0), Cm(10), Cm(10))
    add_text(tf, "매출을 거래 건수와 객단가로 나눈 결과.", 13, BODY, line=1.5, space_after=14, first=True)
    add_text(tf, "9개 중 8개 시장에서 객단가 하락", 13, INK, line=1.5, space_after=10)
    add_text(tf, "동서시장  건수 +330% / 객단가 −44%", 12, BODY, line=1.5, space_after=4)
    add_text(tf, "경동시장  건수 +109% / 객단가 −31%", 12, BODY, line=1.5, space_after=14)
    add_text(tf, "매출 총액은 늘었으나 거래 한 건이 만드는 금액은 줄었다.", 12, MUTED, line=1.5)

    # 6 — 시그니처 (coral)
    signature(prs, CORAL, "발견 01",
              "매출은 늘었지만\n거래 단가는 떨어졌다",
              ["거래 건수는 크게 늘었으나 한 건당 금액은 줄었다. 간편결제로 소액 거래가 기록에 잡힌 영향과 "
               "고액 도매가 소매로 바뀐 영향이 함께 있는 것으로 보인다.",
               "서울약령시장은 건수가 28% 늘고도 객단가가 31% 떨어져 매출이 11.6% 줄었다.",
               "단가를 지탱할 품목이나 서비스가 없으면 물량만으로는 성장을 이어가기 어렵다."],
              "8/9", "객단가가 하락한 시장")

    # 7 — 변이-할당 (white)
    s = blank(prs, CANVAS)
    eyebrow(s, "04  부진 원인 분해")
    title(s, "업종 탓인가 입지 탓인가")
    rule(s, Cm(4.6))
    chart(s, C.diverging([(m, float(r["CE_기여%"])) for m, r in ss.iterrows()],
                         unit="%", label="경쟁력 효과"), "ss", MARGIN, Cm(5.7), Cm(19.4))
    tf = textbox(s, Cm(21.6), Cm(6.0), Cm(10), Cm(10))
    add_text(tf, "매출 변화를 서울 전체 성장분, 업종 구성 효과, 같은 업종 내 상대 성과로 분해했다.",
             13, BODY, line=1.5, space_after=14, first=True)
    add_text(tf, "서울약령시장", 13, INK, line=1.5, space_after=6)
    add_text(tf, "업종 구성 효과  +16.5%", 12, BODY, line=1.5, space_after=3)
    add_text(tf, "경쟁력 효과  −52.9%", 12, CORAL, line=1.5, space_after=14)
    add_text(tf, "보유 업종은 오히려 성장하는 분야였다.", 12, MUTED, line=1.5)

    # 8 — 약령시 업종 (white)
    s = blank(prs, CANVAS)
    eyebrow(s, "04  부진 원인 분해")
    title(s, "모든 주요 업종이 서울 평균에 못 미친다")
    rule(s, Cm(4.6))
    table(s, ["업종", "매출 비중", "서울약령시장", "서울 전통시장", "격차"], [
        ["청과상", "68.0%", "−26.7%", "+52.1%", "−78.8%p"],
        ["의약품", "22.2%", "−0.3%", "+19.1%", "−19.4%p"],
        ["한의원", "4.9%", "−27.4%", "−5.4%", "−22.0%p"],
        ["슈퍼마켓", "2.9%", "−2.1%", "+5.2%", "−7.3%p"],
    ], MARGIN, Cm(5.7), W - MARGIN * 2, [7.0, 5.4, 6.0, 6.0, 5.1], font=12)
    tf = textbox(s, MARGIN, Cm(12.6), Cm(26), Cm(4))
    add_text(tf, "약령시 매출의 68%는 한약재가 아니라 청과상에서 나온다. "
                 "서울 전통시장 청과상이 52% 성장한 기간에 약령시 구역은 27% 줄었다.",
             14, BODY, line=1.5, space_after=8, first=True)
    add_text(tf, "한방 수요 축소는 부진의 일부일 뿐이며, 집객 사업으로는 회복되지 않는다.", 14, INK, line=1.5)

    # 9 — 시그니처 (forest)
    signature(prs, FOREST, "발견 02",
              "약령시의 부진은\n품목이 아니라 입지에서 온다",
              ["업종 구성 효과가 +16.5%라는 것은 보유 업종이 서울 평균보다 잘 성장하는 분야라는 뜻이다.",
               "그럼에도 매출이 339억 원 줄었다. 경쟁력 효과 −52.9%가 업종의 이점을 모두 상쇄했다.",
               "다른 시장과 같은 활성화 처방을 적용해서는 안 된다. 업종 전환과 상권 구조 재편이 필요하다."],
              "−52.9%", "경쟁력 효과")

    # 10 — 배후인구 (white)
    s = blank(prs, CANVAS)
    eyebrow(s, "05  배후 수요")
    title(s, "신규 인구는 시장 도보권 밖에 자리잡았다")
    rule(s, Cm(4.6))
    rows = [[d, f"{ash.loc[d, '0-19_변화율']:+.1f}%", f"{ash.loc[d, '20-39_변화율']:+.1f}%",
             f"{ash.loc[d, '40-59_변화율']:+.1f}%", f"{ash.loc[d, '60+_변화율']:+.1f}%",
             {"전농1동": "롯데캐슬 SKY-L65", "용신동": "한양수자인 그라시엘",
              "청량리동": "청량리종합시장 등", "제기동": "경동시장·서울약령시"}[d]]
            for d in ["전농1동", "용신동", "청량리동", "제기동"]]
    table(s, ["행정동", "0~19세", "20~39세", "40~59세", "60세 이상", "주요 시설"],
          rows, MARGIN, Cm(5.7), W - MARGIN * 2, [4.4, 3.8, 3.8, 3.8, 4.2, 9.5], font=11.5)
    tf = textbox(s, MARGIN, Cm(12.4), Cm(26.5), Cm(4))
    add_text(tf, "2023년 청량리역 일대 약 2,800세대 입주 효과는 전농1동과 용신동에 나타났다. "
                 "경동시장과 서울약령시가 있는 제기동은 같은 기간 인구가 줄었다.",
             14, BODY, line=1.5, space_after=8, first=True)
    add_text(tf, "제기동 감소분은 60대 이상에 집중된다. 시장 매출의 46~72%를 차지하는 연령대다.",
             14, INK, line=1.5)

    # 11 — 인구-매출 검정 (white)
    s = blank(prs, CANVAS)
    eyebrow(s, "06  인과 검정")
    title(s, "인구가 줄면 매출도 주는가")
    rule(s, Cm(4.6))
    table(s, ["설정", "표본", "관측치", "추정 계수", "95% 신뢰구간", "R²"],
          [[r.설정, r.표본, f"{int(r.관측치):,}", f"{r.beta:.2f}",
            f"[{r.ci_lo:.2f}, {r.ci_hi:.2f}]", f"{r.r2:.3f}"] for _, r in el.iterrows()],
          MARGIN, Cm(5.7), W - MARGIN * 2, [6.0, 8.4, 3.6, 3.8, 4.6, 3.1], font=11.5)
    tf = textbox(s, MARGIN, Cm(11.4), Cm(26.5), Cm(5))
    add_text(tf, "청량리 9개만으로는 표본이 작아 서울 전통시장 전체로 확대하고, "
                 "424개 행정동의 야간 생활인구를 상권 경계로 연결해 두 가지 설정으로 추정했다.",
             14, BODY, line=1.5, space_after=10, first=True)
    add_text(tf, "두 설정 모두 계수가 0과 구분되지 않는다. 관측치 5,444개에서 신뢰구간이 "
                 "[−0.15, 0.15]로 좁아 영향이 크다는 가능성까지 배제된다.", 14, INK, line=1.5)

    # 12 — 시그니처 (dark)
    signature(prs, INK, "발견 03",
              "인구가 줄면 매출도 준다\n이 전제는 확인되지 않았다",
              ["제기동 배후인구 감소와 서울약령시장 매출 감소는 각각 관측된 사실이다.",
               "그러나 서울 279개 시장 5,444개 관측치에서 두 변수의 관계는 0과 구분되지 않는다.",
               "인구 지표를 근거로 매출 영향을 추정하거나 지원 규모를 산정하는 방식은 재검토가 필요하다."],
              "0.00", "배후인구-매출 추정 계수")

    # 13 — 시장 유형 (white)
    s = blank(prs, CANVAS)
    eyebrow(s, "07  상권 성격")
    title(s, "한약재 도매는 9개 중 1개 시장뿐이다")
    rule(s, Cm(4.6))
    names = {1: "주말 소매형", 2: "평일 도매형", 3: "혼합 대형", 4: "근린 생필품형"}
    rows = [[m, names[int(r["군집_4"])], f"{dow.loc[m, '주말비중']:.1f}%"]
            for m, r in clus.sort_values("군집_4").iterrows()]
    table(s, ["시장", "거래 패턴 기반 유형", "주말 매출 비중"], rows,
          MARGIN, Cm(5.7), Cm(17.5), [7.0, 6.5, 4.0], font=10.5)
    tf = textbox(s, Cm(21.0), Cm(6.0), Cm(10.6), Cm(11))
    add_text(tf, "시간대 6개, 요일 7개, 연령대 6개로 이루어진 19차원 매출 구성비로 군집화했다.",
             13, BODY, line=1.5, space_after=14, first=True)
    add_text(tf, "서울약령시장만 단독 군집", 13, INK, line=1.5, space_after=6)
    add_text(tf, "주말 비중 19.3%, 일요일 5.4%", 12, BODY, line=1.5, space_after=14)
    add_text(tf, "경동시장은 청과상 35.3%로 종합시장에 가깝고, 청과물·동서시장은 새벽 도매가 아니라 "
                 "주말 소매다.", 12, MUTED, line=1.5)

    # 14 — 공실 통계 (white)
    s = blank(prs, CANVAS)
    eyebrow(s, "08  진단 지표 검증")
    title(s, "공실 통계는 갱신되지 않고 이월된다")
    rule(s, Cm(4.6))
    chart(s, C.hbar([(r["구간"], float(r["동일비율%"])) for _, r in dq.tail(6).iterrows()],
                    unit="%", label="전년 동일 기록 비율",
                    highlight={"2021→2022", "2022→2023"}), "dq", MARGIN, Cm(5.7), Cm(19.4))
    tf = textbox(s, Cm(21.6), Cm(6.0), Cm(10), Cm(10))
    add_text(tf, "전국 1,812개 시장 15개 연도를 검증했다.", 13, BODY, line=1.5, space_after=14, first=True)
    add_text(tf, "최근 구간 47.6%가 전년과 완전히 같은 값", 13, INK, line=1.5, space_after=12)
    add_text(tf, "경동시장은 2020·2021·2022년 3개 연도가 영업 362·빈점포 316으로 동일하다. "
                 "인용되는 46.6%는 관측된 변화가 아니다.", 12, BODY, line=1.5, space_after=12)
    add_text(tf, "원자료 검증 없이는 정반대의 판단이 나올 수 있다.", 12, CORAL, line=1.5)

    # 15 — 실행 과제 (cream)
    s = blank(prs, CREAM)
    eyebrow(s, "09  실행 과제", BODY)
    title(s, "무엇을 먼저 해야 하는가")
    rule(s, Cm(4.6))
    table(s, ["순위", "과제", "근거", "소관", "시점"], [
        ["1", "정비사업 배후인구 설계 검토", "입주분이 시장 도보권 밖에 귀착", "서울시·동대문구", "즉시"],
        ["2", "서울약령시장 개별 진단", "경쟁력 효과 −52.9%", "동대문구·상인회", "즉시"],
        ["3", "성과 지표 개편", "집객이 고객 구성으로 미전환", "사업 시행 주체", "차기 설계"],
        ["4", "고령 고객 접근성 확보", "제기동 60대 이상 −15.1%", "동대문구", "6개월"],
        ["5", "공실 통계 정합성 확인", "전년 동일 기록 47.6%", "소진공", "6개월"],
        ["6", "객단가 하락 원인 규명", "8개 시장 객단가 하락", "분석 기관", "후속"],
    ], MARGIN, Cm(5.7), W - MARGIN * 2, [2.2, 9.0, 9.5, 5.4, 3.4], font=11)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"저장: {OUT}")
    print(f"  슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장 "
          f"· 시그니처 카드 4장 (coral·forest·dark·cream)")


if __name__ == "__main__":
    build()
